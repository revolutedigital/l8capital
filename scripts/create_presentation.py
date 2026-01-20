"""
Script para criar a apresentação PowerPoint da L8 Capital
Requer: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# ============================================
# CONFIGURAÇÕES DE MARCA L8 CAPITAL
# ============================================

# Cores da marca (RGB)
COLORS = {
    'preto_premium': RgbColor(10, 10, 10),      # #0A0A0A
    'branco': RgbColor(255, 255, 255),           # #FFFFFF
    'dourado': RgbColor(201, 162, 39),           # #C9A227
    'dourado_claro': RgbColor(217, 184, 92),     # #D9B85C
    'dourado_escuro': RgbColor(150, 119, 32),    # #967720
    'cinza_600': RgbColor(92, 92, 86),           # #5C5C56
    'cinza_400': RgbColor(168, 168, 160),        # #A8A8A0
    'cinza_200': RgbColor(232, 232, 228),        # #E8E8E4
    'cinza_100': RgbColor(245, 245, 243),        # #F5F5F3
    'cinza_800': RgbColor(58, 58, 54),           # #3A3A36
}

# Fontes (Plus Jakarta Sans pode não estar disponível, usando fallbacks)
FONTS = {
    'display': 'Plus Jakarta Sans',  # Fallback: Segoe UI
    'body': 'Inter',                  # Fallback: Calibri
}

# Dimensões do slide (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_shape_fill(shape, color):
    """Define cor de preenchimento de uma forma"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_shape_no_fill(shape):
    """Remove preenchimento de uma forma"""
    shape.fill.background()


def set_text_frame_properties(text_frame, margin_left=Inches(0), margin_top=Inches(0)):
    """Configura propriedades do text frame"""
    text_frame.word_wrap = True
    text_frame.margin_left = margin_left
    text_frame.margin_top = margin_top
    text_frame.margin_right = margin_left
    text_frame.margin_bottom = margin_top


def add_text(shape, text, font_name='Inter', font_size=18, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Adiciona texto formatado a uma forma"""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_paragraph(text_frame, text, font_name='Inter', font_size=14, bold=False, color=None, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(6)):
    """Adiciona um parágrafo ao text frame"""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def create_slide_with_bg(prs, bg_color):
    """Cria um slide com fundo colorido"""
    slide_layout = prs.slide_layouts[6]  # Layout em branco
    slide = prs.slides.add_slide(slide_layout)

    # Adiciona retângulo de fundo
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(bg, bg_color)
    bg.line.fill.background()

    # Move para trás
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)

    return slide


def create_presentation():
    """Cria a apresentação completa"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ========================================
    # SLIDE 1: CAPA
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['preto_premium'])

    # Logo L8 (texto estilizado)
    logo_box = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(4.333), Inches(1.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "L8"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(120)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # CAPITAL
    capital_box = slide.shapes.add_textbox(Inches(4.5), Inches(3.3), Inches(4.333), Inches(0.6))
    tf = capital_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "C A P I T A L"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(32)
    run.font.bold = False
    run.font.color.rgb = COLORS['dourado']

    # Linha decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.2), Inches(3.333), Pt(2))
    set_shape_fill(line, COLORS['dourado'])
    line.line.fill.background()

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7.333), Inches(0.5))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Sua imobiliária mais forte."
    run.font.name = 'Inter'
    run.font.size = Pt(24)
    run.font.italic = True
    run.font.color.rgb = COLORS['branco']

    # Rodapé
    footer_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(5.333), Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Apresentação Institucional  •  Janeiro 2026"
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 2: NOSSA HISTÓRIA
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Nossa História"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Por que a L8 existe?"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # Quote box
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.733), Inches(1.8))
    set_shape_fill(quote_bg, COLORS['cinza_100'])
    quote_bg.line.fill.background()

    quote_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11), Inches(1.5))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '"Durante 24 anos, trabalhamos do lado das seguradoras. Vimos de perto como as imobiliárias eram tratadas: como mais um canal de vendas, não como parceiras. Comissões baixas, processos burocráticos, zero suporte real. Um dia, decidimos mudar de lado — o lado que importa."'
    run.font.name = 'Inter'
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = COLORS['cinza_800']

    # Timeline
    # Box 1: 2000-2020
    box1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(4.2), Inches(3.2), Inches(2))
    set_shape_fill(box1, COLORS['preto_premium'])
    box1.line.fill.background()

    txt1 = slide.shapes.add_textbox(Inches(1.2), Inches(4.4), Inches(2.8), Inches(1.7))
    tf = txt1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2000-2020"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']
    add_paragraph(tf, "20 anos no mercado de seguros", 'Inter', 14, False, COLORS['branco'], PP_ALIGN.CENTER, Pt(12))
    add_paragraph(tf, "Experiência acumulada", 'Inter', 12, False, COLORS['cinza_400'], PP_ALIGN.CENTER, Pt(6))

    # Seta 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(5), Inches(0.8), Inches(0.4))
    set_shape_fill(arrow1, COLORS['dourado'])
    arrow1.line.fill.background()

    # Box 2: 2021
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(4.2), Inches(3.2), Inches(2))
    set_shape_fill(box2, COLORS['dourado'])
    box2.line.fill.background()

    txt2 = slide.shapes.add_textbox(Inches(5.6), Inches(4.4), Inches(2.8), Inches(1.7))
    tf = txt2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2021"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']
    add_paragraph(tf, "A Decisão", 'Inter', 14, True, COLORS['preto_premium'], PP_ALIGN.CENTER, Pt(12))
    add_paragraph(tf, '"Vamos pro outro lado"', 'Inter', 12, False, COLORS['cinza_800'], PP_ALIGN.CENTER, Pt(6))

    # Seta 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.8), Inches(5), Inches(0.8), Inches(0.4))
    set_shape_fill(arrow2, COLORS['dourado'])
    arrow2.line.fill.background()

    # Box 3: 2022+
    box3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(4.2), Inches(3.2), Inches(2))
    set_shape_fill(box3, COLORS['preto_premium'])
    box3.line.fill.background()

    txt3 = slide.shapes.add_textbox(Inches(10), Inches(4.4), Inches(2.8), Inches(1.7))
    tf = txt3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "2022+"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']
    add_paragraph(tf, "Nasce a L8 Capital", 'Inter', 14, False, COLORS['branco'], PP_ALIGN.CENTER, Pt(12))
    add_paragraph(tf, "Tecnologia própria + Parceria", 'Inter', 12, False, COLORS['cinza_400'], PP_ALIGN.CENTER, Pt(6))

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "02"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 3: NOSSO PROPÓSITO
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Nosso Propósito"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "O Golden Circle da L8"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # Círculo POR QUÊ (centro)
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(2), Inches(3), Inches(3))
    set_shape_fill(circle1, COLORS['dourado'])
    circle1.line.fill.background()

    txt_c1 = slide.shapes.add_textbox(Inches(5.4), Inches(2.8), Inches(2.6), Inches(1.4))
    tf = txt_c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "POR QUÊ?"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']
    add_paragraph(tf, "Imobiliárias merecem parceiro, não fornecedor", 'Inter', 11, False, COLORS['preto_premium'], PP_ALIGN.CENTER, Pt(8))

    # Círculo COMO (meio)
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.2), Inches(1.5), Inches(5), Inches(4))
    circle2.fill.background()
    circle2.line.color.rgb = COLORS['cinza_400']
    circle2.line.width = Pt(2)

    # Círculo O QUÊ (externo)
    circle3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.2), Inches(1), Inches(7), Inches(5))
    circle3.fill.background()
    circle3.line.color.rgb = COLORS['cinza_200']
    circle3.line.width = Pt(2)

    # Textos laterais
    como_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(2.5), Inches(1.5))
    tf = como_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "COMO?"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['cinza_600']
    add_paragraph(tf, "Tecnologia própria + risco compartilhado", 'Inter', 12, False, COLORS['cinza_600'], PP_ALIGN.LEFT, Pt(8))

    oque_box = slide.shapes.add_textbox(Inches(10), Inches(2.5), Inches(2.5), Inches(1.5))
    tf = oque_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "O QUÊ?"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['cinza_600']
    add_paragraph(tf, "Soluções que aumentam receita, reduzem custos e eliminam operacional", 'Inter', 12, False, COLORS['cinza_600'], PP_ALIGN.LEFT, Pt(8))

    # Quote box
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(10.333), Inches(1.2))
    set_shape_fill(quote_bg, COLORS['preto_premium'])
    quote_bg.line.fill.background()

    quote_box = slide.shapes.add_textbox(Inches(1.8), Inches(6), Inches(9.8), Inches(0.9))
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"Quando a imobiliária cresce, todo o ecossistema cresce. Fortalecer imobiliárias é fortalecer o mercado imobiliário brasileiro."'
    run.font.name = 'Inter'
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = COLORS['branco']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "03"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 4: QUEM SOMOS
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['preto_premium'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Quem Somos"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['branco']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "L8 Capital em Números"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_400']

    # Cards de métricas
    metrics = [
        ("24", "ANOS", "Experiência no setor"),
        ("100%", "TECNOLOGIA", "Própria"),
        ("R$ 2.050", "ECONOMIA", "Mensal típica*"),
    ]

    card_width = Inches(3.5)
    start_x = Inches(1.2)
    gap = Inches(0.5)

    for i, (number, label, sublabel) in enumerate(metrics):
        x = start_x + (card_width + gap) * i

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), card_width, Inches(3))
        card.fill.background()
        card.line.color.rgb = COLORS['dourado']
        card.line.width = Pt(2)

        # Número
        num_box = slide.shapes.add_textbox(x, Inches(2.5), card_width, Inches(1.2))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = number
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(56)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        # Label
        lbl_box = slide.shapes.add_textbox(x, Inches(3.7), card_width, Inches(0.5))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['branco']

        # Sublabel
        sub_box = slide.shapes.add_textbox(x, Inches(4.2), card_width, Inches(0.5))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = sublabel
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_400']

    # Nota de rodapé
    note_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.733), Inches(0.3))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "* Base: imobiliária com 500 boletos/mês"
    run.font.name = 'Inter'
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS['cinza_400']

    # Definição
    def_box = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.733), Inches(1))
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Plataforma digital de soluções financeiras para imobiliárias."
    run.font.name = 'Inter'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS['branco']
    add_paragraph(tf, "Tecnologia própria para aumentar receita, reduzir custos e eliminar carga operacional.", 'Inter', 14, False, COLORS['cinza_400'], PP_ALIGN.LEFT, Pt(8))

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "04"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_600']

    # ========================================
    # SLIDE 5: O PROBLEMA
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "O Problema"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "A Dor Real das Imobiliárias em 2026"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # Quote boxes
    quotes = [
        ('"Trabalho 12 horas por dia, mas o lucro não cresce."', Inches(0.8), Inches(1.9)),
        ('"Minha margem diminui todo ano. Os custos só aumentam."', Inches(6.8), Inches(1.9)),
        ('"Todo mundo oferece a mesma coisa. Como me diferencio?"', Inches(0.8), Inches(3.1)),
        ('"Minha equipe gasta mais tempo com papel do que com cliente."', Inches(6.8), Inches(3.1)),
    ]

    for quote, x, y in quotes:
        q_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(1))
        set_shape_fill(q_bg, COLORS['cinza_100'])
        q_bg.line.fill.background()

        q_txt = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(5.4), Inches(0.7))
        tf = q_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = quote
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = COLORS['cinza_800']

    # Tabela de sintomas
    table_data = [
        ("Sintoma", "Causa Raiz", "Impacto"),
        ("Margem Apertada", "Boleto a R$ 6+, comissões baixas", "Lucro estagnado"),
        ("Tempo Perdido", "Processos manuais, retrabalho", "Equipe improdutiva"),
        ("Comoditização", "Mesmos serviços que concorrentes", "Briga só por preço"),
    ]

    table_y = Inches(4.5)
    row_height = Inches(0.6)
    col_widths = [Inches(2.8), Inches(4.5), Inches(3.5)]

    for row_idx, row in enumerate(table_data):
        x = Inches(1.2)
        for col_idx, cell in enumerate(row):
            cell_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x,
                table_y + row_height * row_idx,
                col_widths[col_idx],
                row_height
            )

            if row_idx == 0:
                set_shape_fill(cell_bg, COLORS['preto_premium'])
                text_color = COLORS['branco']
                is_bold = True
            else:
                set_shape_fill(cell_bg, COLORS['cinza_100'] if row_idx % 2 == 0 else COLORS['branco'])
                text_color = COLORS['preto_premium'] if col_idx == 0 else COLORS['cinza_800']
                is_bold = col_idx == 0

            cell_bg.line.color.rgb = COLORS['cinza_200']
            cell_bg.line.width = Pt(0.5)

            cell_txt = slide.shapes.add_textbox(
                x + Inches(0.15),
                table_y + row_height * row_idx + Inches(0.15),
                col_widths[col_idx] - Inches(0.3),
                row_height - Inches(0.1)
            )
            tf = cell_txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell
            run.font.name = 'Inter'
            run.font.size = Pt(12)
            run.font.bold = is_bold
            run.font.color.rgb = text_color

            x += col_widths[col_idx]

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "05"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 6: PROPOSTA DE VALOR
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Nossa Proposta de Valor"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Três Pilares de Transformação"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 3 Pilares
    pillars = [
        ("AUMENTAR\nRECEITA", "Comissões recorrentes + Float rendendo", COLORS['dourado'], COLORS['preto_premium']),
        ("REDUZIR\nCUSTOS", "Boleto de R$ 6 para R$ 1,90", COLORS['preto_premium'], COLORS['branco']),
        ("ELIMINAR\nOPERACIONAL", "Automação de apólices + processos", COLORS['preto_premium'], COLORS['branco']),
    ]

    card_width = Inches(3.6)
    start_x = Inches(0.95)
    gap = Inches(0.4)

    for i, (title, desc, bg_color, text_color) in enumerate(pillars):
        x = start_x + (card_width + gap) * i

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), card_width, Inches(2.8))
        set_shape_fill(card, bg_color)
        card.line.fill.background()

        # Título
        t_box = slide.shapes.add_textbox(x + Inches(0.3), Inches(2.2), card_width - Inches(0.6), Inches(1.2))
        tf = t_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = text_color if i > 0 else COLORS['preto_premium']

        # Descrição
        d_box = slide.shapes.add_textbox(x + Inches(0.3), Inches(3.5), card_width - Inches(0.6), Inches(1))
        tf = d_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_800'] if i == 0 else COLORS['cinza_200']

    # Sinal de +
    for i in range(2):
        x = start_x + card_width + gap / 2 - Inches(0.15) + (card_width + gap) * i
        plus = slide.shapes.add_textbox(x, Inches(2.8), Inches(0.3), Inches(0.5))
        tf = plus.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "+"
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(28)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

    # Seta para baixo
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(4.9), Inches(0.7), Inches(0.6))
    set_shape_fill(arrow, COLORS['dourado'])
    arrow.line.fill.background()

    # Resultado
    result_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.7), Inches(6.333), Inches(0.9))
    set_shape_fill(result_bg, COLORS['dourado'])
    result_bg.line.fill.background()

    result_txt = slide.shapes.add_textbox(Inches(3.7), Inches(5.85), Inches(5.933), Inches(0.6))
    tf = result_txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "IMOBILIÁRIA MAIS FORTE"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Quote
    quote_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.733), Inches(0.4))
    tf = quote_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"Aumenta receita, reduz custo ou elimina operacional? Se não passa por esse filtro, não faz parte da L8."'
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = COLORS['cinza_600']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "06"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 7: PORTFÓLIO DE SOLUÇÕES
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Portfólio de Soluções"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "9 Soluções em 3 Categorias"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # L8 Capital box no topo
    top_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.8), Inches(3.333), Inches(0.7))
    set_shape_fill(top_box, COLORS['preto_premium'])
    top_box.line.fill.background()

    top_txt = slide.shapes.add_textbox(Inches(5.2), Inches(1.95), Inches(2.933), Inches(0.4))
    tf = top_txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "L8 CAPITAL"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # 3 Categorias
    categories = [
        ("SEGUROS", ["Incêndio", "Fiança", "Capitalização"], "+ Faturamento"),
        ("FINANCEIRO", ["Boletagem", "Float", "Fundo de Reserva"], "- Despesas"),
        ("CAPACITAÇÃO", ["Trein. Comercial", "Trein. Operacional", "Vistoria"], "+ Produtividade"),
    ]

    card_width = Inches(3.6)
    start_x = Inches(0.95)
    gap = Inches(0.4)

    for i, (cat_name, items, result) in enumerate(categories):
        x = start_x + (card_width + gap) * i

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.8), card_width, Inches(3.2))
        card.fill.background()
        card.line.color.rgb = COLORS['cinza_200']
        card.line.width = Pt(1)

        # Categoria título
        cat_box = slide.shapes.add_textbox(x + Inches(0.2), Inches(3), card_width - Inches(0.4), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cat_name
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

        # Items
        items_box = slide.shapes.add_textbox(x + Inches(0.3), Inches(3.6), card_width - Inches(0.6), Inches(2))
        tf = items_box.text_frame
        for item in items:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = 'Inter'
            run.font.size = Pt(13)
            run.font.color.rgb = COLORS['cinza_600']
            p.space_after = Pt(4)

        # Resultado
        res_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.5), Inches(6.2), card_width - Inches(1), Inches(0.5))
        set_shape_fill(res_bg, COLORS['dourado'])
        res_bg.line.fill.background()

        res_txt = slide.shapes.add_textbox(x + Inches(0.5), Inches(6.28), card_width - Inches(1), Inches(0.35))
        tf = res_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = result
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "07"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 8: SEGUROS
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Seguros"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Receita Recorrente, Não Venda Única"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # Tabela de produtos
    products = [
        ("Produto", "O que é", "Diferencial L8"),
        ("Seguro Incêndio", "Obrigatório por lei", "Comissão mensal, não única"),
        ("Seguro Fiança", "Garantia locatícia", "Aprovação em 72h"),
        ("Capitalização", "Alternativa ao caução", "100% digital"),
    ]

    table_y = Inches(1.8)
    row_height = Inches(0.65)
    col_widths = [Inches(3), Inches(4), Inches(4.5)]

    for row_idx, row in enumerate(products):
        x = Inches(0.9)
        for col_idx, cell in enumerate(row):
            cell_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x,
                table_y + row_height * row_idx,
                col_widths[col_idx],
                row_height
            )

            if row_idx == 0:
                set_shape_fill(cell_bg, COLORS['preto_premium'])
                text_color = COLORS['branco']
                is_bold = True
            else:
                set_shape_fill(cell_bg, COLORS['branco'])
                text_color = COLORS['dourado'] if col_idx == 2 else COLORS['preto_premium']
                is_bold = col_idx == 0 or col_idx == 2

            cell_bg.line.color.rgb = COLORS['cinza_200']
            cell_bg.line.width = Pt(0.5)

            cell_txt = slide.shapes.add_textbox(
                x + Inches(0.15),
                table_y + row_height * row_idx + Inches(0.18),
                col_widths[col_idx] - Inches(0.3),
                row_height - Inches(0.1)
            )
            tf = cell_txt.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = cell
            run.font.name = 'Inter'
            run.font.size = Pt(13)
            run.font.bold = is_bold
            run.font.color.rgb = text_color

            x += col_widths[col_idx]

    # Comparativo
    comp_title = slide.shapes.add_textbox(Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.4))
    tf = comp_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Comparativo de Ganhos (1 contrato/ano)"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Box tradicional
    trad_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(5.1), Inches(4.5), Inches(1.6))
    trad_box.fill.background()
    trad_box.line.color.rgb = COLORS['cinza_400']
    trad_box.line.width = Pt(1)

    trad_txt = slide.shapes.add_textbox(Inches(1.2), Inches(5.2), Inches(4.1), Inches(1.4))
    tf = trad_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "MODELO TRADICIONAL"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['cinza_600']
    add_paragraph(tf, "Venda única: R$ 200", 'Inter', 12, False, COLORS['cinza_600'], PP_ALIGN.LEFT, Pt(10))
    add_paragraph(tf, "Meses seguintes: R$ 0", 'Inter', 12, False, COLORS['cinza_600'], PP_ALIGN.LEFT, Pt(4))
    add_paragraph(tf, "TOTAL ANO: R$ 200", 'Inter', 14, True, COLORS['cinza_800'], PP_ALIGN.LEFT, Pt(8))

    # Box L8
    l8_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6), Inches(5.1), Inches(4.5), Inches(1.6))
    set_shape_fill(l8_box, COLORS['preto_premium'])
    l8_box.line.fill.background()

    l8_txt = slide.shapes.add_textbox(Inches(6.2), Inches(5.2), Inches(4.1), Inches(1.4))
    tf = l8_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "MODELO L8"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']
    add_paragraph(tf, "Todo mês: R$ 50", 'Inter', 12, False, COLORS['branco'], PP_ALIGN.LEFT, Pt(10))
    add_paragraph(tf, "× 12 meses", 'Inter', 12, False, COLORS['cinza_400'], PP_ALIGN.LEFT, Pt(4))
    add_paragraph(tf, "TOTAL ANO: R$ 600", 'Inter', 14, True, COLORS['dourado'], PP_ALIGN.LEFT, Pt(8))

    # Resultado
    result_txt = slide.shapes.add_textbox(Inches(10.7), Inches(5.5), Inches(2.2), Inches(0.8))
    tf = result_txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "+R$ 400"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']
    add_paragraph(tf, "/contrato/ano", 'Inter', 10, False, COLORS['cinza_600'], PP_ALIGN.CENTER)

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "08"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 9: FINANCEIRO
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['preto_premium'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Financeiro"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['branco']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Economia Real em Boletos"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_400']

    # 3 boxes de comparação
    boxes_data = [
        ("MERCADO", "R$ 6,00", COLORS['cinza_800']),
        ("L8 CAPITAL", "R$ 1,90", COLORS['dourado']),
        ("VOCÊ ECONOMIZA", "R$ 4,10", COLORS['preto_premium']),
    ]

    for i, (label, value, bg_color) in enumerate(boxes_data):
        x = Inches(1.5) + Inches(3.7) * i

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(3.2), Inches(1.6))
        if i == 2:
            set_shape_fill(box, COLORS['dourado'])
        else:
            box.fill.background()
            box.line.color.rgb = COLORS['cinza_600'] if i == 0 else COLORS['dourado']
            box.line.width = Pt(2)

        lbl = slide.shapes.add_textbox(x, Inches(2.1), Inches(3.2), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['cinza_400'] if i < 2 else COLORS['preto_premium']

        val = slide.shapes.add_textbox(x, Inches(2.5), Inches(3.2), Inches(0.8))
        tf = val.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = value
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado'] if i == 1 else (COLORS['preto_premium'] if i == 2 else COLORS['branco'])

        if i == 2:
            sub = slide.shapes.add_textbox(x, Inches(3.2), Inches(3.2), Inches(0.3))
            tf = sub.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "por boleto"
            run.font.name = 'Inter'
            run.font.size = Pt(11)
            run.font.color.rgb = COLORS['preto_premium']

    # Setas
    for i in range(2):
        arrow = slide.shapes.add_textbox(Inches(4.5) + Inches(3.7) * i, Inches(2.5), Inches(0.5), Inches(0.5))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "→"
        run.font.name = 'Arial'
        run.font.size = Pt(32)
        run.font.color.rgb = COLORS['dourado']

    # Tabela de economia
    table_title = slide.shapes.add_textbox(Inches(0.8), Inches(4), Inches(11.733), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Sua Economia Mensal"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS['branco']

    savings_data = [
        ("200 boletos/mês", "× R$ 4,10", "= R$ 820/mês"),
        ("500 boletos/mês", "× R$ 4,10", "= R$ 2.050/mês"),
        ("1.000 boletos/mês", "× R$ 4,10", "= R$ 4.100/mês"),
    ]

    for i, (qty, mult, result) in enumerate(savings_data):
        y = Inches(4.5) + Inches(0.6) * i

        qty_txt = slide.shapes.add_textbox(Inches(2), y, Inches(3), Inches(0.5))
        tf = qty_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = qty
        run.font.name = 'Inter'
        run.font.size = Pt(16)
        run.font.color.rgb = COLORS['branco']

        mult_txt = slide.shapes.add_textbox(Inches(5.5), y, Inches(2), Inches(0.5))
        tf = mult_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = mult
        run.font.name = 'Inter'
        run.font.size = Pt(16)
        run.font.color.rgb = COLORS['cinza_400']

        res_txt = slide.shapes.add_textbox(Inches(8), y, Inches(3), Inches(0.5))
        tf = res_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = result
        run.font.name = 'Inter'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

    # Outras soluções
    other_title = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.4))
    tf = other_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Outras soluções: "
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS['cinza_400']
    run = p.add_run()
    run.text = "Rentabilização de Float"
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']
    run = p.add_run()
    run.text = " (dinheiro parado pode render)  •  "
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS['cinza_400']
    run = p.add_run()
    run.text = "Gestão de Fundo de Reserva"
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "09"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_600']

    # ========================================
    # SLIDE 10: CAPACITAÇÃO
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Capacitação"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Conhecimento que Gera Resultado"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 3 Cards
    cards_data = [
        ("TREINAMENTO COMERCIAL", ["Profissional com experiência real", "Técnicas práticas de conversão", "Customizado para seu mercado"]),
        ("TREINAMENTO OPERACIONAL", ["Metodologias testadas no mercado", "Redução de retrabalho", "Implementação assistida"]),
        ("VISTORIA DE IMÓVEIS", ["Parceiro especializado exclusivo", "Qualidade padronizada", "Preços diferenciados"]),
    ]

    for i, (title, items) in enumerate(cards_data):
        y = Inches(1.85) + Inches(1.8) * i

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), y, Inches(11.333), Inches(1.6))
        card.fill.background()
        card.line.color.rgb = COLORS['cinza_200']
        card.line.width = Pt(1)

        # Título do card
        t_box = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.2), Inches(10.7), Inches(0.4))
        tf = t_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

        # Items
        for j, item in enumerate(items):
            item_box = slide.shapes.add_textbox(Inches(1.3) + Inches(3.8) * j, y + Inches(0.7), Inches(3.5), Inches(0.7))
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = 'Inter'
            run.font.size = Pt(13)
            run.font.color.rgb = COLORS['cinza_600']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "10"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 11: DIFERENCIAL
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Nosso Diferencial"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Parceiro, Não Fornecedor"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # Box Tradicional
    trad_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.8))
    trad_box.fill.background()
    trad_box.line.color.rgb = COLORS['cinza_400']
    trad_box.line.width = Pt(2)

    trad_title = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(5.1), Inches(0.5))
    tf = trad_title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "FORNECEDOR"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS['cinza_600']

    trad_items = [
        "• Vende e vai embora",
        "• Comissão única",
        '• "Precisa de seguro?"',
        "• Tecnologia genérica",
        "• Processo rígido",
        "• Risco: 100% seu",
    ]

    for i, item in enumerate(trad_items):
        item_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.7) + Inches(0.45) * i, Inches(4.9), Inches(0.4))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_600']

    # VS
    vs_box = slide.shapes.add_textbox(Inches(6.3), Inches(3.3), Inches(0.7), Inches(0.6))
    tf = vs_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "VS"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # Box L8
    l8_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.9), Inches(5.5), Inches(3.8))
    set_shape_fill(l8_box, COLORS['preto_premium'])
    l8_box.line.fill.background()

    l8_title = slide.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.5))
    tf = l8_title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PARCEIRO"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    l8_items = [
        "• Cresce junto",
        "• Ganho recorrente",
        '• "Como aumentar sua receita?"',
        "• Tecnologia própria",
        "• Agilidade de startup",
        "• Risco compartilhado",
    ]

    for i, item in enumerate(l8_items):
        item_box = slide.shapes.add_textbox(Inches(7.4), Inches(2.7) + Inches(0.45) * i, Inches(4.9), Inches(0.4))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['branco']

    # Quote
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(6), Inches(9.333), Inches(0.8))
    set_shape_fill(quote_bg, COLORS['dourado'])
    quote_bg.line.fill.background()

    quote_txt = slide.shapes.add_textbox(Inches(2.2), Inches(6.15), Inches(8.933), Inches(0.5))
    tf = quote_txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"24 anos no mercado de seguros. Agora, do lado que importa: O SEU."'
    run.font.name = 'Inter'
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "11"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 12: COMO FUNCIONA
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Como Funciona"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Jornada Simples em 3 Etapas"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 3 Etapas
    steps = [
        ("1", "CONVERSA", "30 min", ["Online ou presencial", "Sem compromisso", "Entender sua operação"]),
        ("2", "DIAGNÓSTICO", "Grátis", ["Análise dos seus custos", "Oportunidades identificadas", "Projeção de economia"]),
        ("3", "IMPLEMENTAÇÃO", "7 dias", ["Integração com sistemas", "Treinamento da equipe", "Suporte contínuo"]),
    ]

    for i, (num, title, badge, items) in enumerate(steps):
        x = Inches(0.8) + Inches(4.2) * i

        # Número
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.3), Inches(1.9), Inches(0.7), Inches(0.7))
        set_shape_fill(num_circle, COLORS['dourado'])
        num_circle.line.fill.background()

        num_txt = slide.shapes.add_textbox(x + Inches(1.3), Inches(2), Inches(0.7), Inches(0.5))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

        # Título
        title_txt = slide.shapes.add_textbox(x, Inches(2.7), Inches(3.5), Inches(0.5))
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

        # Badge
        badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(1), Inches(3.2), Inches(1.5), Inches(0.4))
        set_shape_fill(badge_bg, COLORS['preto_premium'])
        badge_bg.line.fill.background()

        badge_txt = slide.shapes.add_textbox(x + Inches(1), Inches(3.25), Inches(1.5), Inches(0.3))
        tf = badge_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        # Items
        for j, item in enumerate(items):
            item_txt = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.8) + Inches(0.35) * j, Inches(3.3), Inches(0.35))
            tf = item_txt.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"• {item}"
            run.font.name = 'Inter'
            run.font.size = Pt(12)
            run.font.color.rgb = COLORS['cinza_600']

        # Seta (exceto último)
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.6), Inches(2.9), Inches(0.5), Inches(0.3))
            set_shape_fill(arrow, COLORS['dourado'])
            arrow.line.fill.background()

    # Timeline de resultados
    timeline_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.733), Inches(0.4))
    tf = timeline_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Linha do Tempo de Resultados"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    timeline_items = [
        ("MÊS 1", "Economia em boletos no caixa"),
        ("MÊS 2", "Receita de seguros iniciando"),
        ("MÊS 3", "ROI claro e documentado"),
    ]

    for i, (month, desc) in enumerate(timeline_items):
        x = Inches(1.2) + Inches(4) * i

        month_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.7), Inches(3.5), Inches(1))
        set_shape_fill(month_bg, COLORS['cinza_100'])
        month_bg.line.fill.background()

        month_txt = slide.shapes.add_textbox(x + Inches(0.2), Inches(5.85), Inches(3.1), Inches(0.3))
        tf = month_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = month
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        desc_txt = slide.shapes.add_textbox(x + Inches(0.2), Inches(6.15), Inches(3.1), Inches(0.4))
        tf = desc_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(11)
        run.font.color.rgb = COLORS['cinza_600']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "12"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 13: IDENTIDADE VISUAL
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['preto_premium'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Identidade Visual"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['branco']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "A Marca L8 Capital"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_400']

    # Seção Paleta de Cores
    palette_title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5), Inches(0.4))
    tf = palette_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "PALETA DE CORES"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # Cores
    colors_display = [
        ("Preto Premium", "#0A0A0A", COLORS['preto_premium']),
        ("Branco", "#FFFFFF", COLORS['branco']),
        ("Dourado Premium", "#C9A227", COLORS['dourado']),
    ]

    for i, (name, hex_code, color) in enumerate(colors_display):
        x = Inches(0.8) + Inches(3.8) * i

        # Color swatch
        swatch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(3.3), Inches(1.2))
        set_shape_fill(swatch, color)
        if i == 1:  # Branco precisa de borda
            swatch.line.color.rgb = COLORS['cinza_600']
            swatch.line.width = Pt(1)
        else:
            swatch.line.fill.background()

        # Nome da cor
        name_txt = slide.shapes.add_textbox(x, Inches(3.6), Inches(3.3), Inches(0.3))
        tf = name_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = name
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLORS['branco']

        # Código hex
        hex_txt = slide.shapes.add_textbox(x, Inches(3.9), Inches(3.3), Inches(0.3))
        tf = hex_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = hex_code
        run.font.name = 'Inter'
        run.font.size = Pt(11)
        run.font.color.rgb = COLORS['cinza_400']

    # Seção Tipografia
    typo_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(5), Inches(0.4))
    tf = typo_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "TIPOGRAFIA"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # Display font
    display_txt = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(5.5), Inches(1))
    tf = display_txt.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Display: Plus Jakarta Sans"
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS['cinza_400']
    add_paragraph(tf, "AUMENTAR RECEITA", 'Plus Jakarta Sans', 24, True, COLORS['branco'], PP_ALIGN.LEFT, Pt(8))

    # Body font
    body_txt = slide.shapes.add_textbox(Inches(7), Inches(5), Inches(5.5), Inches(1))
    tf = body_txt.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Body: Inter"
    run.font.name = 'Inter'
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS['cinza_400']
    add_paragraph(tf, "Texto corrido com alta legibilidade", 'Inter', 16, False, COLORS['branco'], PP_ALIGN.LEFT, Pt(8))

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "13"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_600']

    # ========================================
    # SLIDE 14: TOM DE VOZ
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Tom de Voz"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Como a L8 se Comunica"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 4 Princípios
    principles = [
        ("DIRETO", '"De R$ 6 para R$ 1,90 por boleto."'),
        ("CONCRETO", '"24 anos de experiência no setor."'),
        ("HONESTO", '"Funciona melhor para +50 imóveis."'),
        ("PARCEIRO", '"Você decide."'),
    ]

    for i, (principle, example) in enumerate(principles):
        y = Inches(1.8) + Inches(0.9) * i

        # Princípio
        p_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(2), Inches(0.7))
        set_shape_fill(p_bg, COLORS['preto_premium'])
        p_bg.line.fill.background()

        p_txt = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.15), Inches(2), Inches(0.4))
        tf = p_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = principle
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        # Exemplo
        e_txt = slide.shapes.add_textbox(Inches(3), y + Inches(0.15), Inches(9.5), Inches(0.5))
        tf = e_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = example
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_600']

    # Frases de Marca
    frases_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.733), Inches(0.4))
    tf = frases_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "FRASES DE MARCA"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    frases = [
        '"Sua imobiliária mais forte."',
        '"Parceiro, não fornecedor."',
        '"Você só tem custo se tiver ganho primeiro."',
    ]

    for i, frase in enumerate(frases):
        f_txt = slide.shapes.add_textbox(Inches(0.8) + Inches(4.2) * i, Inches(6), Inches(4), Inches(0.5))
        tf = f_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = frase
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.italic = True
        run.font.color.rgb = COLORS['dourado']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "14"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 15: COMPLIANCE
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Compliance e Segurança"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Compromisso com Proteção de Dados"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 4 Cards de compliance
    compliance_items = [
        ("LGPD", "COMPLIANT"),
        ("SERVIDORES", "NACIONAIS"),
        ("CRIPTOGRAFIA", "PONTA A PONTA"),
        ("POLÍTICA", "TRANSPARENTE"),
    ]

    for i, (title, subtitle) in enumerate(compliance_items):
        x = Inches(0.8) + Inches(3.15) * i

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.9), Inches(1.5))
        set_shape_fill(card, COLORS['preto_premium'])
        card.line.fill.background()

        t_txt = slide.shapes.add_textbox(x, Inches(2.2), Inches(2.9), Inches(0.4))
        tf = t_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        s_txt = slide.shapes.add_textbox(x, Inches(2.65), Inches(2.9), Inches(0.4))
        tf = s_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['branco']

        # Checkmark
        check = slide.shapes.add_textbox(x + Inches(1.15), Inches(2.9), Inches(0.6), Inches(0.4))
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "✓"
        run.font.name = 'Arial'
        run.font.size = Pt(20)
        run.font.color.rgb = COLORS['dourado']

    # Lista de garantias
    guarantees = [
        "Criptografia de ponta a ponta em todas as transações",
        "Servidores seguros em território nacional",
        "Políticas rígidas de controle de acesso",
        "Informações disponíveis sob demanda",
    ]

    for i, guarantee in enumerate(guarantees):
        g_txt = slide.shapes.add_textbox(Inches(1.5), Inches(3.8) + Inches(0.45) * i, Inches(10.5), Inches(0.4))
        tf = g_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"•  {guarantee}"
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_600']

    # Quote
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.8), Inches(10.333), Inches(0.9))
    set_shape_fill(quote_bg, COLORS['cinza_100'])
    quote_bg.line.fill.background()

    quote_txt = slide.shapes.add_textbox(Inches(1.8), Inches(6), Inches(9.733), Inches(0.5))
    tf = quote_txt.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"Seus dados e os dados dos seus clientes são tratados com o mais alto padrão de segurança."'
    run.font.name = 'Inter'
    run.font.size = Pt(14)
    run.font.italic = True
    run.font.color.rgb = COLORS['cinza_800']

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "15"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 16: PRÓXIMOS PASSOS
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['branco'])

    # Header dourado
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15))
    set_shape_fill(header, COLORS['dourado'])
    header.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Próximos Passos"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(0.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Como Começar"
    run.font.name = 'Inter'
    run.font.size = Pt(20)
    run.font.color.rgb = COLORS['cinza_600']

    # 3 Passos
    steps_data = [
        ("1", "AGENDAR", "30 min sem compromisso"),
        ("2", "RECEBER", "Diagnóstico grátis"),
        ("3", "DECIDIR", "Sem pressão"),
    ]

    for i, (num, title, desc) in enumerate(steps_data):
        x = Inches(0.8) + Inches(4.2) * i

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(2))
        set_shape_fill(card, COLORS['preto_premium'])
        card.line.fill.background()

        # Número
        num_txt = slide.shapes.add_textbox(x, Inches(2), Inches(3.7), Inches(0.8))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = COLORS['dourado']

        # Título
        t_txt = slide.shapes.add_textbox(x, Inches(2.8), Inches(3.7), Inches(0.4))
        tf = t_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title
        run.font.name = 'Plus Jakarta Sans'
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = COLORS['branco']

        # Descrição
        d_txt = slide.shapes.add_textbox(x, Inches(3.2), Inches(3.7), Inches(0.4))
        tf = d_txt.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = desc
        run.font.name = 'Inter'
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS['cinza_400']

        # Setas
        if i < 2:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(3.8), Inches(2.6), Inches(0.35), Inches(0.25))
            set_shape_fill(arrow, COLORS['dourado'])
            arrow.line.fill.background()

    # Contato
    contact_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.733), Inches(0.4))
    tf = contact_title.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "CANAIS DE CONTATO"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']

    # Contatos
    contacts = [
        ("Email:", "contato@l8capital.com.br"),
        ("Site:", "l8capital.com.br"),
        ("LinkedIn:", "linkedin.com/company/l8capital"),
    ]

    for i, (label, value) in enumerate(contacts):
        c_txt = slide.shapes.add_textbox(Inches(0.8), Inches(4.7) + Inches(0.4) * i, Inches(6), Inches(0.35))
        tf = c_txt.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{label}  "
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS['cinza_600']
        run = p.add_run()
        run.text = value
        run.font.name = 'Inter'
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS['preto_premium']

    # WhatsApp QR placeholder
    qr_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(4.5), Inches(4.5), Inches(2))
    set_shape_fill(qr_bg, COLORS['cinza_100'])
    qr_bg.line.fill.background()

    qr_txt = slide.shapes.add_textbox(Inches(8.2), Inches(4.7), Inches(4.1), Inches(1.6))
    tf = qr_txt.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "WhatsApp"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS['preto_premium']
    add_paragraph(tf, "[QR Code]", 'Inter', 32, False, COLORS['cinza_400'], PP_ALIGN.CENTER, Pt(10))
    add_paragraph(tf, "Resposta em 24h", 'Inter', 11, False, COLORS['cinza_600'], PP_ALIGN.CENTER, Pt(8))

    # Número do slide
    num_box = slide.shapes.add_textbox(Inches(12.5), Inches(7), Inches(0.5), Inches(0.3))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "16"
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_400']

    # ========================================
    # SLIDE 17: ENCERRAMENTO
    # ========================================
    slide = create_slide_with_bg(prs, COLORS['preto_premium'])

    # Logo L8 (texto estilizado)
    logo_box = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(4.333), Inches(1.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "L8"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(120)
    run.font.bold = True
    run.font.color.rgb = COLORS['dourado']

    # CAPITAL
    capital_box = slide.shapes.add_textbox(Inches(4.5), Inches(3.3), Inches(4.333), Inches(0.6))
    tf = capital_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "C A P I T A L"
    run.font.name = 'Plus Jakarta Sans'
    run.font.size = Pt(32)
    run.font.bold = False
    run.font.color.rgb = COLORS['dourado']

    # Linha decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.2), Inches(3.333), Pt(2))
    set_shape_fill(line, COLORS['dourado'])
    line.line.fill.background()

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(3), Inches(4.5), Inches(7.333), Inches(0.5))
    tf = tagline_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Sua imobiliária mais forte."
    run.font.name = 'Inter'
    run.font.size = Pt(24)
    run.font.italic = True
    run.font.color.rgb = COLORS['branco']

    # Contatos
    contact_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.333), Inches(0.8))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "contato@l8capital.com.br"
    run.font.name = 'Inter'
    run.font.size = Pt(14)
    run.font.color.rgb = COLORS['cinza_400']
    add_paragraph(tf, "l8capital.com.br", 'Inter', 14, False, COLORS['cinza_400'], PP_ALIGN.CENTER, Pt(4))

    # Copyright
    copyright_box = slide.shapes.add_textbox(Inches(3), Inches(6.6), Inches(7.333), Inches(0.4))
    tf = copyright_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "© 2026 L8 Capital. Todos os direitos reservados."
    run.font.name = 'Inter'
    run.font.size = Pt(10)
    run.font.color.rgb = COLORS['cinza_600']

    # ========================================
    # SALVAR APRESENTAÇÃO
    # ========================================
    output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'L8_Capital_Apresentacao.pptx')
    prs.save(output_path)
    print(f"Apresentação salva em: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
